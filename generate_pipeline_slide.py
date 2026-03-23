#!/usr/bin/env python3
"""
Generate a single PowerPoint slide showing the LiDAR-Camera fusion +
segmentation pipeline for the AVMI UGV research paper.
"""

import pptx
import pptx.enum.shapes
import pptx.enum.text
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import pptx.oxml.ns as nsmap
from lxml import etree

# ── helpers ────────────────────────────────────────────────────────────────────

def rgb(r, g, b):
    return RGBColor(r, g, b)

def add_box(slide, left, top, width, height,
            text, fill_color, text_color=None,
            font_size=11, bold=False, border_color=None, border_width=Pt(1.5),
            line2=None, font_size2=9):
    """Add a rounded-rectangle box with centred text."""
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 5,  # ROUNDED_RECTANGLE = 5
        left, top, width, height
    )
    shape.shape_type  # access to trigger init
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_color

    line = shape.line
    line.width = border_width
    if border_color:
        line.color.rgb = border_color
    else:
        # darken fill by ~40 for border
        line.color.rgb = RGBColor(
            max(0, fill_color[0] - 40),
            max(0, fill_color[1] - 40),
            max(0, fill_color[2] - 40),
        )

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    from pptx.util import Pt as PT
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree

    # vertical centering
    tf.vertical_anchor = pptx.enum.text.MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(font_size)
    run.font.color.rgb = text_color or RGBColor(255, 255, 255)

    if line2:
        from pptx.oxml import parse_xml
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = line2
        r2.font.size = Pt(font_size2)
        r2.font.italic = True
        r2.font.color.rgb = text_color or RGBColor(230, 230, 230)

    return shape


def add_arrow(slide, x1, y1, x2, y2, color=RGBColor(80, 80, 80), width=Pt(2.0)):
    """Add a line with an arrowhead from (x1,y1) to (x2,y2) in EMU."""
    from pptx.util import Emu
    connector = slide.shapes.add_connector(
        pptx.enum.shapes.MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2
    )
    connector.line.width = width
    connector.line.color.rgb = color
    # add end arrowhead via XML
    ln = connector.line._ln
    tail_elem = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd')
    # actually we want head arrow (at the end of the line direction)
    head_elem = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}headEnd')
    head_elem.set('type', 'arrow')
    head_elem.set('w', 'med')
    head_elem.set('len', 'med')
    return connector


def add_label(slide, left, top, width, height, text, color, size=8.5, italic=False, bold=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.bold = bold
    return tb


# ══════════════════════════════════════════════════════════════════════════════
#  BUILD SLIDE
# ══════════════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width  = Inches(16)   # wide (16:9)
prs.slide_height = Inches(9)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

# ── Background ────────────────────────────────────────────────────────────────
bg = slide.background
bg_fill = bg.fill
bg_fill.solid()
bg_fill.fore_color.rgb = RGBColor(245, 247, 250)  # very light grey-blue

# ── Title bar ─────────────────────────────────────────────────────────────────
title_box = slide.shapes.add_shape(5,
    Inches(0), Inches(0), Inches(16), Inches(0.7))
title_box.fill.solid()
title_box.fill.fore_color.rgb = RGBColor(30, 50, 80)
title_box.line.fill.background()

tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.08), Inches(15.4), Inches(0.55))
tf = tb.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Semantic LiDAR–Camera Fusion Pipeline for UGV Terrain Perception"
r.font.size = Pt(18)
r.font.bold = True
r.font.color.rgb = RGBColor(255, 255, 255)

# ── Column X positions (centre of each column) ───────────────────────────────
# Layout (inches):
#  Col0=Input sensors  Col1=Processing  Col2=Fusion  Col3=Output
#
#  [Camera]   [LiDAR]
#      ↓          ↓
#  [GANav Seg] [Coord Transform]
#      ↓          ↓
#  [Seg Mask]  [Projected Pts]
#       \         /
#        [FUSION ]
#        /       \
#  [Image Overlay] [Semantic PCD]

C = [1.3, 4.7, 8.1, 11.5, 14.7]   # centre X
W  = Inches(2.2)                    # box width
H  = Inches(0.72)                   # box height standard

def cx(col): return Inches(C[col]) - W/2   # left edge from centre
def cy(row): return Inches(row)             # top edge

# ────────────────────────────────────────────────────────────────────────────
#  ROW 0 — Section labels
# ────────────────────────────────────────────────────────────────────────────
for lbl, xi in [('① Data Acquisition', 1.3),
                ('② Pre-processing', 5.9),
                ('③ Semantic Fusion', 9.6),
                ('④ Outputs', 13.2)]:
    add_label(slide,
              Inches(xi - 1.3), Inches(0.78), Inches(2.6), Inches(0.28),
              lbl, RGBColor(30, 50, 80), size=9, bold=True)

# ────────────────────────────────────────────────────────────────────────────
#  ROW 1 — Input sensors
# ────────────────────────────────────────────────────────────────────────────
# Camera box
CAM_L, CAM_T = Inches(0.25), Inches(1.15)
add_box(slide, CAM_L, CAM_T, W, H,
        "RGB Camera", rgb(42, 110, 160),
        font_size=11, bold=True,
        line2="Raw image frames")

# LiDAR box
LID_L, LID_T = Inches(2.65), Inches(1.15)
add_box(slide, LID_L, LID_T, W, H,
        "Semantic LiDAR", rgb(160, 82, 30),
        font_size=11, bold=True,
        line2="UE packed RGB labels")

# Separator line between input and processing columns
sep = slide.shapes.add_shape(1,  # rectangle
    Inches(4.95), Inches(0.72), Pt(1.5), Inches(8.0))
sep.fill.solid(); sep.fill.fore_color.rgb = RGBColor(180, 190, 210)
sep.line.fill.background()

sep2 = slide.shapes.add_shape(1,
    Inches(8.75), Inches(0.72), Pt(1.5), Inches(8.0))
sep2.fill.solid(); sep2.fill.fore_color.rgb = RGBColor(180, 190, 210)
sep2.line.fill.background()

sep3 = slide.shapes.add_shape(1,
    Inches(12.25), Inches(0.72), Pt(1.5), Inches(8.0))
sep3.fill.solid(); sep3.fill.fore_color.rgb = RGBColor(180, 190, 210)
sep3.line.fill.background()

# ────────────────────────────────────────────────────────────────────────────
#  ROW 2 — Processing branch 1: Segmentation  (left sub-column ~x=5.2)
#           Processing branch 2: LiDAR xform  (right sub-column ~x=8.0)
# ────────────────────────────────────────────────────────────────────────────
SEG_L,  SEG_T  = Inches(5.10), Inches(1.15)
XFRM_L, XFRM_T = Inches(7.45), Inches(1.15)

add_box(slide, SEG_L, SEG_T, W, H,
        "GANav Model", rgb(46, 125, 50),
        font_size=11, bold=True,
        line2="MiT-B0 + PSA decode head")

add_box(slide, XFRM_L, XFRM_T, W, H,
        "Coordinate Transform", rgb(123, 70, 155),
        font_size=11, bold=True,
        line2="P_cam = R·P_lidar + t")

# ────────────────────────────────────────────────────────────────────────────
#  ROW 3 — Outputs of processing
# ────────────────────────────────────────────────────────────────────────────
SMASK_L, SMASK_T = Inches(5.10), Inches(2.45)
PROJ_L,  PROJ_T  = Inches(7.45), Inches(2.45)

add_box(slide, SMASK_L, SMASK_T, W, H,
        "Segmentation Mask", rgb(56, 142, 60),
        font_size=11, bold=True,
        line2="6-class pixel labels")

add_box(slide, PROJ_L, PROJ_T, W, H,
        "Projected Points", rgb(142, 86, 180),
        font_size=11, bold=True,
        line2="u = fx·(X/Z)+cx")

# ────────────────────────────────────────────────────────────────────────────
#  ROW 4 — Pinhole + depth filter note (small annotation box)
# ────────────────────────────────────────────────────────────────────────────
add_box(slide, Inches(7.45), Inches(3.40), W, Inches(0.58),
        "Depth Filter", rgb(170, 110, 200),
        font_size=9.5, bold=False,
        line2="0.5 m < Z < 50 m", font_size2=8.5)

# ────────────────────────────────────────────────────────────────────────────
#  FUSION BOX (centre)
# ────────────────────────────────────────────────────────────────────────────
FUS_W = Inches(2.8)
FUS_H = Inches(1.1)
FUS_L = Inches(8.95) - FUS_W/2
FUS_T = Inches(4.35)

add_box(slide, FUS_L, FUS_T, FUS_W, FUS_H,
        "Semantic Fusion", rgb(183, 28, 28),
        font_size=13, bold=True,
        line2="Semantic colour assignment\nto projected 3-D points",
        font_size2=9, border_width=Pt(2.5))

# ────────────────────────────────────────────────────────────────────────────
#  OUTPUT COLUMN — Image Overlay
# ────────────────────────────────────────────────────────────────────────────
OV_L,  OV_T  = Inches(12.45), Inches(2.90)
PCD_L, PCD_T = Inches(12.45), Inches(4.70)
OUT_W = Inches(3.0)
OUT_H = Inches(0.95)

add_box(slide, OV_L, OV_T, OUT_W, OUT_H,
        "Semantic Overlay Image", rgb(13, 71, 161),
        font_size=11, bold=True,
        line2="Camera frame with coloured LiDAR dots")

add_box(slide, PCD_L, PCD_T, OUT_W, OUT_H,
        "Semantic Point Cloud", rgb(0, 96, 100),
        font_size=11, bold=True,
        line2="3-D cloud coloured by terrain class")

# Class legend box
LEG_L, LEG_T = Inches(12.45), Inches(6.05)
add_box(slide, LEG_L, LEG_T, OUT_W, Inches(1.7),
        "Terrain Classes", rgb(50, 60, 75),
        font_size=10, bold=True,
        border_width=Pt(1.2))

# Legend items (coloured small boxes)
legend = [
    ("Sky",        rgb(100, 160, 220)),
    ("Tree",       rgb(34, 139, 34)),
    ("Bush",       rgb(107, 142, 35)),
    ("Ground",     rgb(139, 100, 43)),
    ("Obstacle",   rgb(160, 160, 60)),
    ("Rock",       rgb(180, 60, 60)),
]
for i, (lbl, col) in enumerate(legend):
    lx = LEG_L + Inches(0.18) + (i % 3) * Inches(0.95)
    ly = LEG_T + Inches(0.55) + (i // 3) * Inches(0.52)
    add_box(slide, lx, ly, Inches(0.80), Inches(0.35),
            lbl, col, font_size=8, bold=False, border_width=Pt(0.8))

# ────────────────────────────────────────────────────────────────────────────
#  UE SOURCE note
# ────────────────────────────────────────────────────────────────────────────
add_box(slide, Inches(2.65), Inches(2.50), W, Inches(0.58),
        "UE Semantic Labels", rgb(195, 110, 40),
        font_size=9.5, bold=False,
        line2="Rock / Ground / Vegetation", font_size2=8.5)

# ────────────────────────────────────────────────────────────────────────────
#  ARROWS
# ────────────────────────────────────────────────────────────────────────────
A = RGBColor(60, 60, 60)

def mid_right(l, t, w, h): return (l + w, t + h/2)
def mid_left(l, t, w, h):  return (l,     t + h/2)
def mid_bot(l, t, w, h):   return (l + w/2, t + h)
def mid_top(l, t, w, h):   return (l + w/2, t)

# Camera → GANav
add_arrow(slide, *mid_right(CAM_L, CAM_T, W, H), *mid_left(SEG_L, SEG_T, W, H), A)
# LiDAR → Coord Transform
add_arrow(slide, *mid_right(LID_L, LID_T, W, H), *mid_left(XFRM_L, XFRM_T, W, H), A)
# GANav → Seg Mask
add_arrow(slide, *mid_bot(SEG_L, SEG_T, W, H), *mid_top(SMASK_L, SMASK_T, W, H), A)
# Coord Transform → Projected Points
add_arrow(slide, *mid_bot(XFRM_L, XFRM_T, W, H), *mid_top(PROJ_L, PROJ_T, W, H), A)
# Projected Points → Depth Filter
add_arrow(slide, *mid_bot(PROJ_L, PROJ_T, W, H),
          Inches(7.45) + W/2, Inches(3.40), A)
# UE Labels → LiDAR (small annotation)
add_arrow(slide, *mid_top(LID_L + Inches(0.0), Inches(2.50), W, Inches(0.58)),
          *mid_bot(LID_L, LID_T, W, H), A)

# Seg Mask → Fusion
sx, sy = mid_bot(SMASK_L, SMASK_T, W, H)
fx, fy = mid_top(FUS_L, FUS_T, FUS_W, FUS_H)
add_arrow(slide, sx, sy, fx - Inches(0.35), fy, A)

# Depth Filter → Fusion
px, py = Inches(7.45) + W/2, Inches(3.40) + Inches(0.58)
add_arrow(slide, px, py, fx + Inches(0.35), fy, A)

# Fusion → Image Overlay
fx2, fy2 = mid_right(FUS_L, FUS_T, FUS_W, FUS_H)
add_arrow(slide, fx2, fy2 - Inches(0.25),
          *mid_left(OV_L, OV_T, OUT_W, OUT_H), A)

# Fusion → Semantic PCD
add_arrow(slide, fx2, fy2 + Inches(0.25),
          *mid_left(PCD_L, PCD_T, OUT_W, OUT_H), A)

# ────────────────────────────────────────────────────────────────────────────
#  Inline annotation labels on arrows
# ────────────────────────────────────────────────────────────────────────────
def ann(slide, x, y, text):
    tb = slide.shapes.add_textbox(Inches(x) - Inches(0.7), Inches(y) - Inches(0.14),
                                  Inches(1.4), Inches(0.30))
    tf = tb.text_frame
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r  = p.add_run()
    r.text = text
    r.font.size = Pt(7.5)
    r.font.italic = True
    r.font.color.rgb = RGBColor(90, 90, 90)

ann(slide, 3.3, 1.50, "image frames")
ann(slide, 6.5, 1.50, "UE bag data")
ann(slide, 9.55, 3.90, "filtered pts (u,v,Z)")
ann(slide, 9.10, 5.90, "2-D overlay")
ann(slide, 10.50, 5.15, "3-D coloured")

# ────────────────────────────────────────────────────────────────────────────
#  Footer
# ────────────────────────────────────────────────────────────────────────────
foot = slide.shapes.add_shape(1,
    Inches(0), Inches(8.70), Inches(16), Inches(0.30))
foot.fill.solid(); foot.fill.fore_color.rgb = RGBColor(30, 50, 80)
foot.line.fill.background()

tb = slide.shapes.add_textbox(Inches(0.3), Inches(8.72), Inches(15), Inches(0.25))
tf = tb.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
r = p.add_run()
r.text = ("AVMI UGV Dataset  |  Sensor fusion: Unreal Engine calibration  |  "
          "Segmentation: GANav (MiT-B0 + PSA)  |  Output: 6-class semantic overlay + coloured point cloud")
r.font.size = Pt(7.5)
r.font.color.rgb = RGBColor(200, 210, 230)

# ── Save ─────────────────────────────────────────────────────────────────────
out = '/home/pinaka/GANav-offroad/pipeline_slide.pptx'
prs.save(out)
print(f"Saved: {out}")
